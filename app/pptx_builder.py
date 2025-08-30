# app/pptx_builder.py
"""
Builds a PowerPoint (.pptx) from an outline and an uploaded template.

Upgrades:
- Optional exact image reuse with safe placement (avoid overlapping text, add images before text).
- Capability-based layout fallback (ensure Title + Body/Content placeholders exist).
- Template slides cleared safely (drop rels) after harvesting pictures (prevents "repair" prompts).
- Theme fonts/colors applied consistently.
- Two-column split for dense slide bullets; sub-bullets ("  • ") rendered at level-1.

This module is defensive against unusual templates and missing placeholders.
"""

from __future__ import annotations

from io import BytesIO
from typing import List, Optional, Tuple, Dict

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

from .schemas import Outline, OutlineSlide
from .template_utils import (
    extract_template_images,   # fallback media scraping
    find_preferred_layout,     # name-based helper
    get_theme_style,           # theme colors & fonts {"colors": {...}, "fonts": {...}}
)

# ---------------- Tunables ----------------

TITLE_FALLBACK_SIZE_PT = 32
BODY_FALLBACK_SIZE_PT = 18
MAX_BULLETS_PER_SLIDE = 12
EMU_PER_INCH = 914400

# Recognize sub-bullets (the parser uses a simple "  • " prefix for nesting).
SUB_BULLET_PREFIX = "  • "

# ---------------- Low-level helpers ----------------

def _strip_control_chars(s: str) -> str:
    # Remove control chars except \t \n
    return "".join(ch for ch in (s or "") if ch == "\t" or ch == "\n" or ord(ch) >= 32)

def _rgb_from_hex(hex6: Optional[str]) -> Optional[RGBColor]:
    if not hex6 or len(hex6) < 6:
        return None
    try:
        r = int(hex6[0:2], 16)
        g = int(hex6[2:4], 16)
        b = int(hex6[4:6], 16)
        return RGBColor(r, g, b)
    except Exception:
        return None

def _apply_font_to_runs(text_frame, *, name: Optional[str], size_pt: Optional[int], color: Optional[RGBColor]):
    # Apply font to all runs across all paragraphs in a text_frame.
    for p in text_frame.paragraphs:
        for r in p.runs:
            if name:
                r.font.name = name
            if size_pt:
                r.font.size = Pt(size_pt)
            if color:
                r.font.color.rgb = color

def _bullet_level_and_text(text: str) -> Tuple[int, str]:
    """
    Detect whether this bullet should be level-0 or level-1.
    The parser formats sub-bullets with a '  • ' prefix; detect this and set level=1.
    """
    s = _strip_control_chars(text or "").rstrip()
    if not s:
        return 0, ""
    if s.startswith(SUB_BULLET_PREFIX):
        return 1, s[len(SUB_BULLET_PREFIX):].strip()
    # Also tolerate literal "•"
    if s.strip().startswith("•") and not s.startswith("•"):  # e.g., "   • text"
        cleaned = s[s.index("•") + 1:].strip()
        return 1, cleaned
    if s.startswith("• "):
        return 1, s[2:].strip()
    return 0, s

# ---------------- Placeholder finders ----------------

def _title_placeholder(slide):
    for shp in slide.placeholders:
        try:
            if shp.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                return shp
        except Exception:
            continue
    return slide.placeholders[0] if slide.placeholders else None

def _subtitle_placeholder(slide):
    for shp in slide.placeholders:
        try:
            if shp.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:
                return shp
        except Exception:
            continue
    return None

def _content_placeholders(slide):
    """Text-capable content placeholders (BODY/CONTENT/SUBTITLE but not TITLE)."""
    out = []
    for shp in slide.placeholders:
        try:
            ptype = shp.placeholder_format.type
            if ptype in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.CONTENT, PP_PLACEHOLDER.SUBTITLE):
                _ = shp.text_frame  # ensure text-capable
                out.append(shp)
        except Exception:
            continue
    return out

def _picture_placeholders(slide):
    holders = []
    for shp in slide.placeholders:
        try:
            if shp.placeholder_format.type in (PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.CONTENT):
                holders.append(shp)
        except Exception:
            continue
    return holders

def _first_picture_placeholder(slide):
    pics = _picture_placeholders(slide)
    return pics[0] if pics else None

# ---------------- Geometry & safe-zone helpers ----------------

def _rect(left: int, top: int, width: int, height: int) -> Dict[str, int]:
    return {"left": max(0, left), "top": max(0, top), "width": max(0, width), "height": max(0, height)}

def _collect_text_zones(slide) -> List[Dict[str, int]]:
    """
    Rects for any shape likely to hold text: title/body/center-title/subtitle/content placeholders,
    plus any autoshape/textbox with a text_frame.
    """
    zones: List[Dict[str, int]] = []
    for sh in slide.shapes:
        try:
            if getattr(sh, "is_placeholder", False):
                ph_type = getattr(getattr(sh, "placeholder_format", None), "type", None)
                # common text placeholders: 1=TITLE, 2=BODY, 3=CENTER_TITLE, 4=SUBTITLE, 7=CONTENT
                if ph_type in (1, 2, 3, 4, 7):
                    zones.append(_rect(int(sh.left), int(sh.top), int(sh.width), int(sh.height)))
                    continue
            if getattr(sh, "has_text_frame", False):
                zones.append(_rect(int(sh.left), int(sh.top), int(sh.width), int(sh.height)))
        except Exception:
            continue
    return zones

def _intersect_area(a: Dict[str, int], b: Dict[str, int]) -> int:
    if not a or not b:
        return 0
    ax1, ay1, ax2, ay2 = a["left"], a["top"], a["left"] + a["width"], a["top"] + a["height"]
    bx1, by1, bx2, by2 = b["left"], b["top"], b["left"] + b["width"], b["top"] + b["height"]
    ix1, iy1, ix2, iy2 = max(ax1, bx1), max(ay1, by1), min(ax2, bx2), min(ay2, by2)
    if ix2 <= ix1 or iy2 <= iy1:
        return 0
    return (ix2 - ix1) * (iy2 - iy1)

def _overlaps_any_text(img: Dict[str, int], zones: List[Dict[str, int]], thresh: float = 0.10) -> bool:
    area = max(1, img["width"] * img["height"])
    for z in zones:
        if _intersect_area(img, z) / area > thresh:
            return True
    return False

def _choose_safe_zone(slide_w: int, slide_h: int,
                      title_rect: Optional[Dict[str, int]],
                      body_rect: Optional[Dict[str, int]],
                      pad: int = int(0.1 * EMU_PER_INCH)) -> Dict[str, int]:
    """
    Prefer a column to the RIGHT of the body; if too narrow, use BELOW the body.
    If no body placeholder, fall back to area under the title; else a right sidebar.
    """
    if body_rect:
        # Right of body
        right_left = body_rect["left"] + body_rect["width"] + pad
        right_width = max(0, slide_w - right_left - pad)
        right_top = body_rect["top"]
        right_height = body_rect["height"]
        if right_width >= slide_w * 0.18 and right_height >= slide_h * 0.18:
            return _rect(right_left, right_top, right_width, right_height)

        # Below body
        below_top = body_rect["top"] + body_rect["height"] + pad
        below_height = max(0, slide_h - below_top - pad)
        if below_height >= slide_h * 0.18:
            return _rect(pad, below_top, max(0, slide_w - 2 * pad), below_height)

        # Left of body (last resort)
        left_width = max(0, body_rect["left"] - 2 * pad)
        if left_width >= slide_w * 0.18:
            return _rect(pad, body_rect["top"], left_width, body_rect["height"])

    if title_rect:
        area_top = title_rect["top"] + title_rect["height"] + pad
        area_height = max(0, slide_h - area_top - pad)
        return _rect(pad, area_top, max(0, slide_w - 2 * pad), area_height)

    # Fallback: right sidebar
    sidebar_left = int(slide_w * 0.64) + pad
    sidebar_width = max(0, int(slide_w * 0.36) - 2 * pad)
    sidebar_top = int(slide_h * 0.18) + pad
    sidebar_height = max(0, int(slide_h * 0.72) - 2 * pad)
    return _rect(sidebar_left, sidebar_top, sidebar_width, sidebar_height)

def _fit_into_box(img: Dict[str, int], box: Dict[str, int]) -> Dict[str, int]:
    """Scale img to fit within box, keep aspect ratio; center it inside the box."""
    iw, ih = max(1, img["width"]), max(1, img["height"])
    bw, bh = max(1, box["width"]), max(1, box["height"])
    scale = min(bw / iw, bh / ih, 1.0)
    nw, nh = int(iw * scale), int(ih * scale)
    nl = box["left"] + (bw - nw) // 2
    nt = box["top"] + (bh - nh) // 2
    return _rect(nl, nt, nw, nh)

# ---------------- Template picture harvesting & slide clearing ----------------

def _harvest_template_pictures(prs: Presentation) -> List[List[Dict[str, int]]]:
    """
    Collect per-slide picture specs BEFORE clearing slides, so they can be re-inserted
    into generated slides later. Returns a list indexed by slide order; each item is a list
    of dicts: {"blob": bytes, "left": int, "top": int, "width": int, "height": int}
    """
    collected: List[List[Dict[str, int]]] = []
    for s in prs.slides:
        specs: List[Dict[str, int]] = []
        for shape in s.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    specs.append({
                        "blob": shape.image.blob,
                        "left": int(shape.left), "top": int(shape.top),
                        "width": int(shape.width), "height": int(shape.height),
                    })
            except Exception:
                continue
        collected.append(specs)
    return collected

def _clear_all_slides_safely(prs: Presentation) -> None:
    """
    Delete all slides and drop relationships to prevent 'repair' prompts in Office.
    """
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.rId
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

def _find_title_and_content_layout_index(prs: Presentation) -> Optional[int]:
    """
    Heuristic: find a layout that has both a title and a body/content placeholder.
    This is sturdier than name matching on non-English or custom templates.
    """
    for i, layout in enumerate(prs.slide_layouts):
        has_title, has_body = False, False
        try:
            for ph in layout.placeholders:
                t = ph.placeholder_format.type
                if t == PP_PLACEHOLDER.TITLE:        # 1
                    has_title = True
                if t in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.CONTENT):  # 2 or 7
                    has_body = True
            if has_title and has_body:
                return i
        except Exception:
            continue
    return None

def _find_title_layout_index(prs: Presentation) -> Optional[int]:
    """
    Try to find a good title slide layout (Title Slide / Section Header / Title Only).
    Fall back to any layout with a Title placeholder.
    """
    # First try common names
    named = find_preferred_layout(prs, ["Title Slide", "Section Header", "Title Only", "Title"])
    if named is not None:
        for idx, layout in enumerate(prs.slide_layouts):
            if layout is named:
                return idx
    # Then capability search: any layout with a title placeholder
    for i, layout in enumerate(prs.slide_layouts):
        try:
            for ph in layout.placeholders:
                if ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    return i
        except Exception:
            continue
    return 0 if prs.slide_layouts else None

# ---------------- Writers ----------------

def _set_title(slide, title_text: str, theme: Optional[dict]):
    ph = _title_placeholder(slide)
    if ph is None:
        return
    ph.text = _strip_control_chars(title_text or "")

    # Title styling: prefer major font + accent1 color; fallback to dk1
    font_name = (theme or {}).get("fonts", {}).get("major") or (theme or {}).get("fonts", {}).get("minor")
    color_hex = (theme or {}).get("colors", {}).get("accent1") or (theme or {}).get("colors", {}).get("dk1")
    color = _rgb_from_hex(color_hex)
    tf = ph.text_frame
    tf.word_wrap = True
    _apply_font_to_runs(tf, name=font_name, size_pt=TITLE_FALLBACK_SIZE_PT, color=color)

def _set_subtitle_if_present(slide, subtitle_text: Optional[str], theme: Optional[dict]):
    if not subtitle_text:
        return
    ph = _subtitle_placeholder(slide)
    if ph is None:
        return
    ph.text = _strip_control_chars(subtitle_text)
    font_name = (theme or {}).get("fonts", {}).get("minor") or (theme or {}).get("fonts", {}).get("major")
    color_hex = (theme or {}).get("colors", {}).get("dk1") or (theme or {}).get("colors", {}).get("lt1")
    color = _rgb_from_hex(color_hex)
    tf = ph.text_frame
    tf.word_wrap = True
    _apply_font_to_runs(tf, name=font_name, size_pt=BODY_FALLBACK_SIZE_PT, color=color)

def _set_bullets_single(tf, bullets: List[str], theme: Optional[dict]):
    """Write bullets into a single text frame, respecting sub-bullet levels."""
    tf.clear()
    tf.word_wrap = True
    if not bullets:
        return

    bullets = list(bullets)[:MAX_BULLETS_PER_SLIDE]

    # First bullet
    lvl, txt = _bullet_level_and_text(bullets[0])
    p0 = tf.paragraphs[0]
    p0.level = max(0, min(4, lvl))
    p0.text = txt

    # Others
    for b in bullets[1:]:
        lvl, txt = _bullet_level_and_text(b)
        if not txt:
            continue
        para = tf.add_paragraph()
        para.level = max(0, min(4, lvl))
        para.text = txt

    # Apply theme
    body_font = (theme or {}).get("fonts", {}).get("minor") or (theme or {}).get("fonts", {}).get("major")
    body_color_hex = (theme or {}).get("colors", {}).get("dk1")
    body_color = _rgb_from_hex(body_color_hex)
    _apply_font_to_runs(tf, name=body_font, size_pt=BODY_FALLBACK_SIZE_PT, color=body_color)

def _set_bullets(slide, bullets: List[str], theme: Optional[dict]):
    """
    Fill bullets into available content placeholders.
    - If two or more content placeholders exist and there are many bullets, split across first two.
    - Otherwise, write into the first content placeholder.
    """
    placeholders = _content_placeholders(slide)
    if not placeholders:
        return

    bullets = list(bullets)[:MAX_BULLETS_PER_SLIDE]

    # Decide split for "Two Content" layouts when two text placeholders are present.
    if len(placeholders) >= 2 and len(bullets) >= 6:
        half = (len(bullets) + 1) // 2
        left, right = bullets[:half], bullets[half:]
        _set_bullets_single(placeholders[0].text_frame, left, theme)
        _set_bullets_single(placeholders[1].text_frame, right, theme)
        return

    # Single placeholder
    _set_bullets_single(placeholders[0].text_frame, bullets, theme)

# ---------------- Public Builder ----------------

def build_presentation(
    outline: Outline,
    template_bytes: bytes,
    *,
    subtitle: Optional[str] = None,
    reuse_images: bool = False,
) -> bytes:
    """
    Build a .pptx as bytes from the given outline and template (.pptx/.potx bytes).

    Behavior:
    - Harvests template pictures per slide (if reuse_images=True), then clears slides safely.
    - Adds a title slide and content slides using name-based or capability-based layout detection.
    - Inserts images BEFORE text; avoids covering text by relocating/resizing into a safe zone.
    - Applies theme fonts/colors for titles and bullets.
    - Supports two-column bullets when layout provides two content placeholders.
    - Writes speaker notes if present.
    """
    prs = Presentation(BytesIO(template_bytes))
    theme = get_theme_style(template_bytes) or {"colors": {}, "fonts": {}}

    # --- Harvest images per slide (exact reuse) BEFORE clearing slides
    template_pictures: List[List[Dict[str, int]]] = []
    if reuse_images:
        template_pictures = _harvest_template_pictures(prs)

    # --- Clear existing slides to prevent "repair" messages while keeping masters/themes
    _clear_all_slides_safely(prs)

    # --- Choose title layout
    title_layout_idx = _find_title_layout_index(prs)
    if title_layout_idx is None:
        title_layout_idx = 0

    # ----- Title slide -----
    title_slide = prs.slides.add_slide(prs.slide_layouts[title_layout_idx])
    _set_title(title_slide, outline.title, theme)
    _set_subtitle_if_present(title_slide, subtitle, theme)

    # --- Choose a default content layout (Title + Content) by capability
    content_layout_idx = _find_title_and_content_layout_index(prs)
    if content_layout_idx is None:
        # Fall back to a name-based search
        named = find_preferred_layout(prs, ["Title and Content", "Two Content", "Content with Caption", "Blank"])
        if named is not None:
            for idx, layout in enumerate(prs.slide_layouts):
                if layout is named:
                    content_layout_idx = idx
                    break
    if content_layout_idx is None:
        content_layout_idx = 1 if len(prs.slide_layouts) > 1 else 0

    slide_w, slide_h = int(prs.slide_width), int(prs.slide_height)

    # ----- Content slides -----
    for idx, s in enumerate(outline.slides):
        requested = (s.layout or "auto").strip().lower()

        # Select layout: try requested name; else use capability-based default.
        chosen_layout = None
        if requested != "auto":
            chosen_layout = find_preferred_layout(
                prs,
                [s.layout, "Title and Content", "Two Content", "Content with Caption", "Picture with Caption", "Blank"],
            )
        if chosen_layout is None:
            chosen_layout = prs.slide_layouts[content_layout_idx]

        slide = prs.slides.add_slide(chosen_layout)

        # Compute likely text zones (from placeholders) BEFORE placing images
        title_rect = None
        body_rect = None
        if slide.shapes.title:
            t = slide.shapes.title
            title_rect = _rect(int(t.left), int(t.top), int(t.width), int(t.height))
        for ph in _content_placeholders(slide):
            body_rect = _rect(int(ph.left), int(ph.top), int(ph.width), int(ph.height))
            break  # first content area is enough

        text_zones = _collect_text_zones(slide)

        # --- 1) Insert images FIRST so text stays on top; avoid overlapping text
        if reuse_images and template_pictures and idx < len(template_pictures):
            for pic in template_pictures[idx]:
                target = _rect(pic["left"], pic["top"], pic["width"], pic["height"])
                if _overlaps_any_text(target, text_zones, thresh=0.10):
                    safe = _choose_safe_zone(
                        slide_w=slide_w, slide_h=slide_h,
                        title_rect=title_rect, body_rect=body_rect,
                        pad=int(0.1 * EMU_PER_INCH),
                    )
                    target = _fit_into_box(target, safe)
                try:
                    slide.shapes.add_picture(
                        BytesIO(pic["blob"]),
                        target["left"], target["top"],
                        width=target["width"], height=target["height"]
                    )
                except Exception:
                    # Non-fatal: continue placing others
                    pass
        else:
            # Fallback: opportunistic media reuse from /ppt/media if layout suggests picture
            layout_name = (getattr(chosen_layout, "name", "") or "").lower()
            wants_picture = "picture" in layout_name or _first_picture_placeholder(slide) is not None
            media = extract_template_images(template_bytes) if wants_picture else []
            if media:
                # Place a single picture to add some visual continuity
                try:
                    ph = _first_picture_placeholder(slide)
                    if ph is not None:
                        ph.insert_picture(BytesIO(media[0]))
                    else:
                        # Right-side placement
                        width_in = 3.0
                        left_in = max(0.0, (slide_w / EMU_PER_INCH) - width_in - 0.5)
                        top_in = 1.0
                        slide.shapes.add_picture(BytesIO(media[0]), Inches(left_in), Inches(top_in), width=Inches(width_in))
                except Exception:
                    pass

        # --- 2) Now add text so it stays above images
        _set_title(slide, s.title, theme)
        _set_bullets(slide, list(s.bullets or []), theme)

        # Speaker notes (optional)
        if getattr(s, "notes", None) is not None:
            try:
                notes_slide = slide.notes_slide
                notes_tf = notes_slide.notes_text_frame
                notes_tf.clear()
                notes_tf.text = _strip_control_chars(s.notes or "")
            except Exception:
                pass

    # ----- Save -----
    bio = BytesIO()
    prs.save(bio)
    return bio.getvalue()

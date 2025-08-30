# app/template_utils.py
"""
Utilities for working with uploaded PowerPoint templates:
- Safe image extraction from ppt/media (raster formats only)
- Theme parsing (colors + major/minor fonts) from ppt/theme/theme*.xml
- Robust layout selection by name and placeholder capabilities
- PPTX safety checks to avoid zip-bombs / corrupt files
- Slide dimension helper (EMU, inches, cm)
- Template analyzer for debugging (names, capabilities, theme, dimensions, image count)
"""

from __future__ import annotations

import hashlib
import os
import zipfile
from io import BytesIO
from typing import Dict, List, Optional, Tuple
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.slide import SlideLayout
from pptx.enum.shapes import PP_PLACEHOLDER

# ---------------- Limits / Env ----------------

MAX_TEMPLATE_IMAGES = int(os.getenv("MAX_TEMPLATE_IMAGES", "20"))
MAX_TEMPLATE_IMAGE_MB = int(os.getenv("MAX_TEMPLATE_IMAGE_MB", "5"))

# Zip safety
MAX_ZIP_ENTRIES = int(os.getenv("MAX_ZIP_ENTRIES", "2000"))
MAX_ZIP_MEMBER_MB = int(os.getenv("MAX_ZIP_MEMBER_MB", "50"))
MAX_ZIP_TOTAL_MB = int(os.getenv("MAX_ZIP_TOTAL_MB", "200"))
MAX_COMPRESSION_RATIO = float(os.getenv("MAX_COMPRESSION_RATIO", "200.0"))

RASTER_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff"}

# Unit constants
EMU_PER_INCH = 914400
CM_PER_INCH = 2.54

# Build a tolerant placeholder-name map based on what's available in this python-pptx version
_PLACEHOLDER_NAMES: Dict[int, str] = {}
for _name in ("TITLE", "BODY", "CENTER_TITLE", "SUBTITLE", "DATE", "SLIDE_NUMBER", "FOOTER", "HEADER", "CONTENT", "PICTURE"):
    _val = getattr(PP_PLACEHOLDER, _name, None)
    if _val is not None:  # only include if present in this install
        _PLACEHOLDER_NAMES[int(_val)] = _name

# ---------------- Safety ----------------

def is_safe_pptx(
    template_bytes: bytes,
    max_entries: int = MAX_ZIP_ENTRIES,
    max_member_mb: int = MAX_ZIP_MEMBER_MB,
    max_total_mb: int = MAX_ZIP_TOTAL_MB,
    max_ratio: float = MAX_COMPRESSION_RATIO,
) -> bool:
    from zipfile import ZipFile, BadZipFile

    try:
        with ZipFile(BytesIO(template_bytes)) as z:
            names = z.namelist()
            if "[Content_Types].xml" not in names:
                return False
            if not any(n.endswith("ppt/presentation.xml") for n in names):
                return False
            if len(names) > max_entries:
                return False

            per_limit = max_member_mb * 1024 * 1024
            total_uncompressed = 0

            for info in z.infolist():
                if info.is_dir():
                    continue
                if info.file_size > per_limit:
                    return False
                total_uncompressed += info.file_size
                if info.compress_size > 0:
                    ratio = float(info.file_size) / float(info.compress_size)
                    if ratio > max_ratio:
                        return False

            if total_uncompressed > max_total_mb * 1024 * 1024:
                return False

        return True
    except BadZipFile:
        return False
    except Exception:
        return False

# ---------------- Images ----------------

def extract_template_images(template_bytes: bytes) -> List[bytes]:
    images: List[bytes] = []
    seen_hashes = set()
    per_image_limit = MAX_TEMPLATE_IMAGE_MB * 1024 * 1024

    with zipfile.ZipFile(BytesIO(template_bytes)) as z:
        for name in sorted(z.namelist()):
            if not name.startswith("ppt/media/"):
                continue
            ext = os.path.splitext(name)[1].lower()
            if ext not in RASTER_EXTS:
                continue
            with z.open(name) as f:
                data = f.read()
                if len(data) > per_image_limit:
                    continue
                h = hashlib.sha256(data).hexdigest()
                if h in seen_hashes:
                    continue
                images.append(data)
                seen_hashes.add(h)
                if len(images) >= MAX_TEMPLATE_IMAGES:
                    break
    return images

# ---------------- Theme parsing ----------------

def get_theme_style(template_bytes: bytes) -> Dict[str, Dict[str, str]]:
    try:
        with zipfile.ZipFile(BytesIO(template_bytes)) as z:
            theme_name = next((n for n in z.namelist() if n.startswith("ppt/theme/theme")), None)
            if not theme_name:
                return {"colors": {}, "fonts": {}}
            xml_bytes = z.read(theme_name)
    except Exception:
        return {"colors": {}, "fonts": {}}

    try:
        root = ET.fromstring(xml_bytes)
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

        colors: Dict[str, str] = {}
        for tag in ["dk1", "lt1", "dk2", "lt2", "accent1", "accent2", "accent3", "accent4", "accent5", "accent6"]:
            node = root.find(f".//a:clrScheme/a:{tag}", ns)
            if node is None:
                continue
            srgb = node.find(".//a:srgbClr", ns)
            if srgb is not None and "val" in srgb.attrib:
                colors[tag] = srgb.attrib["val"]

        fonts: Dict[str, str] = {}
        major = root.find(".//a:fontScheme/a:majorFont/a:latin", ns)
        minor = root.find(".//a:fontScheme/a:minorFont/a:latin", ns)
        if major is not None and "typeface" in major.attrib:
            fonts["major"] = major.attrib["typeface"]
        if minor is not None and "typeface" in minor.attrib:
            fonts["minor"] = minor.attrib["typeface"]

        return {"colors": colors, "fonts": fonts}
    except Exception:
        return {"colors": {}, "fonts": {}}

# ---------------- Layout selection ----------------

def _layout_capabilities(layout: SlideLayout) -> Tuple[int, int]:
    text_capable = 0
    picture_capable = 0
    try:
        for shp in layout.placeholders:
            try:
                ptype = shp.placeholder_format.type
            except Exception:
                continue
            try:
                _ = shp.text_frame
                text_capable += 1
            except Exception:
                pass
            if ptype in (getattr(PP_PLACEHOLDER, "PICTURE", -1), getattr(PP_PLACEHOLDER, "CONTENT", -1)):
                picture_capable += 1
    except Exception:
        pass
    return text_capable, picture_capable

def _name_match_score(candidate: str, target: str) -> int:
    c = (candidate or "").strip().lower()
    t = (target or "").strip().lower()
    if not c or not t:
        return 0
    if c == t:
        return 3
    if t in c:
        return 2
    return 0

def _capability_ok(target: str, text_count: int, pic_count: int) -> bool:
    t = (target or "").lower()
    if "two content" in t:
        return text_count >= 2
    if "picture" in t:
        return pic_count >= 1
    if "content" in t or "caption" in t:
        return text_count >= 1
    return True

def find_preferred_layout(prs: Presentation, preferred_names: List[str]) -> Optional[SlideLayout]:
    if not preferred_names:
        return None

    lowered = [p.lower() for p in preferred_names]
    for layout in prs.slide_layouts:
        try:
            if layout.name and layout.name.lower() in lowered:
                txt, pic = _layout_capabilities(layout)
                target = preferred_names[lowered.index(layout.name.lower())]
                if _capability_ok(target, txt, pic):
                    return layout
        except Exception:
            continue

    best: Tuple[int, Optional[SlideLayout]] = (0, None)
    for layout in prs.slide_layouts:
        lname = getattr(layout, "name", "") or ""
        txt, pic = _layout_capabilities(layout)
        score = 0
        for pref in preferred_names:
            score = max(score, _name_match_score(lname, pref))
            if score > 0 and not _capability_ok(pref, txt, pic):
                score = 0
        if score > best[0]:
            best = (score, layout)

    return best[1] if best[0] > 0 else None

# ---------------- Dimensions ----------------

def get_ppt_dimensions(prs: Presentation) -> Dict[str, float]:
    w_emu = int(prs.slide_width)
    h_emu = int(prs.slide_height)
    w_in = w_emu / EMU_PER_INCH
    h_in = h_emu / EMU_PER_INCH
    w_cm = w_in * CM_PER_INCH
    h_cm = h_in * CM_PER_INCH
    return {
        "width_emu": w_emu,
        "height_emu": h_emu,
        "width_in": w_in,
        "height_in": h_in,
        "width_cm": w_cm,
        "height_cm": h_cm,
    }

# ---------------- Template Analyzer (for /api/template_info) ----------------

def _placeholder_summary(layout: SlideLayout) -> Dict[str, int | List[str]]:
    names: List[str] = []
    text_capable = 0
    picture_capable = 0
    try:
        for shp in layout.placeholders:
            try:
                ptype = shp.placeholder_format.type
                friendly = _PLACEHOLDER_NAMES.get(int(ptype), f"#{int(ptype)}")
                names.append(friendly)
                try:
                    _ = shp.text_frame
                    text_capable += 1
                except Exception:
                    pass
                if int(ptype) in (
                    int(getattr(PP_PLACEHOLDER, "PICTURE", -1)),
                    int(getattr(PP_PLACEHOLDER, "CONTENT", -1)),
                ):
                    picture_capable += 1
            except Exception:
                continue
    except Exception:
        pass
    return {
        "count": len(names),
        "text_capable": text_capable,
        "picture_capable": picture_capable,
        "types": names,
    }

def analyze_template(template_bytes: bytes) -> Dict[str, object]:
    prs = Presentation(BytesIO(template_bytes))

    dims = get_ppt_dimensions(prs)
    theme = get_theme_style(template_bytes)
    images = extract_template_images(template_bytes)

    layouts: List[Dict[str, object]] = []
    for i, layout in enumerate(prs.slide_layouts):
        try:
            lname = getattr(layout, "name", "") or f"Layout {i}"
        except Exception:
            lname = f"Layout {i}"
        layouts.append({
            "index": i,
            "name": lname,
            "placeholders": _placeholder_summary(layout),
        })

    masters: List[str] = []
    try:
        for m in prs.slide_masters:
            masters.append(getattr(m, "name", "") or "Master")
    except Exception:
        pass

    return {
        "dimensions": dims,
        "theme": theme,
        "image_count": len(images),
        "layouts": layouts,
        "masters": masters,
    }
